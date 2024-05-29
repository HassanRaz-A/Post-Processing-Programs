import os
import tkinter as tk
from tkinter import filedialog, messagebox
from tkinter.scrolledtext import ScrolledText
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from tkinter import ttk
from tkinter.font import Font

# Function to process PowerPoint files
def align_and_adjust_pictures(file_paths, weight, log_text):
    for pptx_path in file_paths:
        log_text.insert(tk.END, f"Processing {pptx_path}...\n")
        log_text.update()

        try:
            prs = Presentation(pptx_path)
        except Exception as e:
            log_text.insert(tk.END, f"Error opening the presentation {pptx_path}: {e}\n")
            log_text.update()
            continue

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    left = int((prs.slide_width - shape.width) / 2)
                    top = int((prs.slide_height - shape.height) / 2)
                    shape.left = left
                    shape.top = top

                    picture_format = shape.line
                    if picture_format is not None:
                        picture_format.width = Pt(float(weight))
                    else:
                        picture_format = shape.line.format
                        picture_format.width = Pt(float(weight))

        try:
            prs.save(pptx_path)
            log_text.insert(tk.END, f"Pictures aligned to center and border weights adjusted successfully for {pptx_path}\n")
            log_text.update()
        except Exception as e:
            log_text.insert(tk.END, f"Error saving the presentation {pptx_path}: {e}\n")
            log_text.update()
    messagebox.showinfo("Processing Complete", "All PowerPoint files processed successfully!")
# Function to browse files and trigger processing
def browse_files(file_entry, weight_combo, log_text):
    file_paths = filedialog.askopenfilenames(filetypes=[("PowerPoint files", "*.pptx")])
    if file_paths:
        file_entry.delete(0, tk.END)
        file_entry.insert(tk.END, "\n".join(file_paths))

# Create Tkinter window
root = tk.Tk()
root.title("PowerPoint Picture Alignment Tool")
root.geometry("700x400")
root.configure(bg='#e1d0ba')  # Set background color

# Bold font
bold_font = Font(family="Helvetica", size=12, weight="bold")

# File entry label and widget
file_label = tk.Label(root, bg='#e1d0ba', text="Select PowerPoint file(s):", font=('Helvetica Bold', 12))
file_label.grid(row=0, column=0, sticky="w", padx=10, pady=10)  # Align to the west (left)
file_entry = tk.Entry(root,bg='#EEEEEE', width=60)
file_entry.grid(row=0, column=1, padx=10, sticky="we")  # Align to the west and east (left and right)

# Browse button for files
browse_button = tk.Button(root, bg='#e1d0ba', text="Browse", font=('Helvetica Bold', 12), command=lambda: browse_files(file_entry, weight_combo, log_text))
browse_button.grid(row=0, column=2, padx=10)  # Set position and size

# Weight selection label and widget
weight_values = [0.5, 1, 1.5, 2, 2.5, 3, 3.5, 4, 4.5, 5, 5.5, 6, 6.5, 7, 7.5, 8, 8.5, 9, 9.5, 10]
# weight_label = tk.Label(root, bg='#e1d0ba', text="Weight:", font=('Helvetica Bold', 12))
# weight_label.grid(row=1, column=1, sticky="ne")  # Align to the west (left)
weight_combo = ttk.Combobox(root, values=weight_values, width=10, state="readonly")
weight_combo.current(0)  # Set default selection
weight_combo.grid(row=1, column=2,sticky="ns", pady=10)  # Align to the west and east (left and right)

# Log display label and widget
log_label = tk.Label(root, bg='#e1d0ba', text="Log:", font=('Helvetica Bold', 12))
log_label.grid(row=1, column=0, sticky="w",padx=10)  # Align to the west (left)
log_text = ScrolledText(root,bg='#EEEEEE', height=12, width=80)
log_text.grid(row=2, column=0, columnspan=3, padx=10 ,sticky="wns")  # Span across three columns

# Execute button
execute_button = tk.Button(root, bg='#e1d0ba', text="Execute", font=('Helvetica Bold', 12),command=lambda: align_and_adjust_pictures(file_entry.get().split("\n"), weight_combo.get(), log_text))
execute_button.grid(row=4, column=0, columnspan=3, padx=10, pady=20, sticky="we")  # Span across three columns

# Start GUI event loop
root.mainloop()

