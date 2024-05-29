import os
import tkinter as tk
from tkinter import filedialog, messagebox
from tkinter.scrolledtext import ScrolledText
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Pt
from tkinter import ttk
from tkinter.font import Font
import tkinter.simpledialog as sd
# Function to adjust the border width of pictures in PowerPoint files
from pptx.dml.color import RGBColor

def adjust_picture_border(file_paths, weight, log_text):
    # Black color (RGB)
    black_color = RGBColor(0, 0, 0)

    for pptx_path in file_paths:
        log_text.insert(tk.END, f"Processing {pptx_path}...\n")
        log_text.update()

        try:
            prs = Presentation(pptx_path)
        except Exception as e:
            log_text.insert(tk.END, f"Error opening the presentation {pptx_path}: {e}\n")
            log_text.update()
            continue

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    # Create a border for the picture shape
                    picture_format = shape.line
                    if picture_format is not None:
                        picture_format.fill.solid()
                        picture_format.fill.fore_color.rgb = black_color
                        picture_format.width = Pt(float(weight))
                    else:
                        picture_format = shape.line.format
                        picture_format.fill.solid()
                        picture_format.fill.fore_color.rgb = black_color
                        picture_format.width = Pt(float(weight))

        try:
            prs.save(pptx_path)
            log_text.insert(tk.END, f"Border width adjusted successfully for {pptx_path}\n")
            log_text.update()
        except Exception as e:
            log_text.insert(tk.END, f"Error saving the presentation {pptx_path}: {e}\n")
            log_text.update()
    messagebox.showinfo("Processing Complete", "All PowerPoint files processed successfully!")



# Function to browse files and trigger processing
def browse_files(file_entry, weight_combo, log_text):
    file_paths = []
    file_type = [("PowerPoint files", "*.pptx")]

    # Create a themed dialog window
    dialog = tk.Toplevel()
    dialog.title("File Selection")
    dialog.geometry("300x180")
    dialog.configure(bg='#87CEEB')  # Set background color to sky blue

    # Create a style for buttons with larger size and hover effects
    style = ttk.Style()
    style.configure("Custom.TButton", font=('Helvetica Bold', 12), background='#87CEEB', padding=10)
    style.map("Custom.TButton", background=[("active", "#add8e6")])  # Hover effect

    # Label and buttons for the options
    label = tk.Label(dialog, text="Select an option:", font=('Helvetica Bold', 12), bg='#87CEEB')
    label.pack(pady=10)

    individual_button = ttk.Button(dialog, text="Individual Files", style="Custom.TButton", font=('Helvetica Bold', 12), command=lambda: select_files(dialog, file_entry, file_type))
    individual_button.pack(pady=5)

    folder_button = ttk.Button(dialog, text="Folder", font=('Helvetica Bold', 12), style="Custom.TButton", command=lambda: select_folder(dialog, file_entry))
    folder_button.pack(pady=5)

def select_files(dialog, file_entry, file_type):
    # Close the dialog
    dialog.destroy()
    # User wants to select individual files
    file_paths = filedialog.askopenfilenames(filetypes=file_type)
    if file_paths:
        file_entry.delete(0, tk.END)
        file_entry.insert(tk.END, "\n".join(file_paths))

def select_folder(dialog, file_entry):
    # Close the dialog
    dialog.destroy()
    # User wants to select a folder
    folder_path = filedialog.askdirectory()
    if folder_path:
        # Get all PowerPoint files in the selected folder
        file_paths = [os.path.join(folder_path, file_name) for file_name in os.listdir(folder_path) if file_name.endswith('.pptx')]
        if file_paths:
            file_entry.delete(0, tk.END)
            file_entry.insert(tk.END, "\n".join(file_paths))



# Create Tkinter window
root = tk.Tk()
root.title("W8 Adjustment Tool")
root.geometry("700x400")
root.configure(bg='#e1d0ba')  # Set background color

# Bold font
bold_font = Font(family="Helvetica", size=10, weight="bold")

# File entry label and widget
file_label = tk.Label(root, bg='#e1d0ba', text="Select PowerPoint file(s):", font=('Helvetica Bold', 12))
file_label.grid(row=0, column=0, sticky="w", padx=10, pady=10)  # Align to the west (left)
file_entry = tk.Entry(root,bg='#EEEEEE', width=60)
file_entry.grid(row=0, column=1, padx=10, pady=10, sticky="we")  # Align to the west and east (left and right)

# Browse button for files
browse_button = tk.Button(root, bg='#e1d0ba', text="Browse", font=('Helvetica Bold', 12), command=lambda: browse_files(file_entry, weight_combo, log_text))
browse_button.grid(row=0, column=2, padx=10, pady=10)  # Set position and size

# Weight selection label and widget
weight_values = [1.5,1,2,2.5,3,3.5,4,4.5,6,6.5]  # Updated weight values
weight_label = tk.Label(root, bg='#e1d0ba', text="Weight:", font=('Helvetica Bold', 12))
weight_label.grid(row=1, column=0, sticky="w", padx=10, pady=10)  # Align to the west (left)
weight_combo = ttk.Combobox(root, values=weight_values, width=5, state="readonly")
weight_combo.current(0)  # Set default selection
weight_combo.grid(row=1, column=1, padx=10, pady=10, sticky="we")  # Align to the west and east (left and right)

# Log display label and widget
log_label = tk.Label(root, bg='#e1d0ba', text="Log:", font=('Helvetica Bold', 12))
log_label.grid(row=2, column=0, sticky="w", padx=10, pady=10)  # Align to the west (left)
log_text = ScrolledText(root,  bg='#EEEEEE',height=10, width=60)
log_text.grid(row=3, column=0, columnspan=3, padx=10, pady=10, sticky="we")  # Span across three columns

# Execute button
execute_button = tk.Button(root, bg='#e1d0ba', text="Execute", font=('Helvetica Bold', 12), command=lambda: adjust_picture_border(file_entry.get().split("\n"), str(weight_combo.get()), log_text))
execute_button.grid(row=4, column=0, columnspan=3, padx=10, pady=10, sticky="we")  # Span across three columns

# Start GUI event loop
root.mainloop()
